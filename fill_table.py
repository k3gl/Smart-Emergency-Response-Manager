# -*- coding: utf-8 -*-
"""Resize Table 1 to 5x5 and fill it with the model-performance results."""

import copy
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn

FILE = "ERM_Poster_FILLED.pptx"

DATA = [
    ["Model / Task",        "Accuracy", "Precision", "Recall", "F1-score"],
    ["Severity (4-class)",  "90.6%",    "0.906",     "0.906",  "0.906"],
    ["Fire needed",         "98.6%",    "0.980",     "0.973",  "0.976"],
    ["Ambulance needed",    "96.4%",    "0.972",     "0.944",  "0.958"],
    ["Police needed",       "95.6%",    "0.962",     "0.968",  "0.965"],
]

prs = Presentation(FILE)
tbl_shape = next(sh for sh in prs.slides[0].shapes if sh.has_table)
tbl = tbl_shape.table._tbl

cur_rows = len(tbl_shape.table.rows)
cur_cols = len(tbl_shape.table.columns)
need_rows, need_cols = len(DATA), len(DATA[0])

# --- add columns (copy last gridCol + a tc into every row) ---
grid = tbl.find(qn("a:tblGrid"))
for _ in range(need_cols - cur_cols):
    grid.append(copy.deepcopy(grid.findall(qn("a:gridCol"))[-1]))
    for tr in tbl.findall(qn("a:tr")):
        tr.append(copy.deepcopy(tr.findall(qn("a:tc"))[-1]))

# --- add rows (copy last tr) ---
for _ in range(need_rows - cur_rows):
    tbl.append(copy.deepcopy(tbl.findall(qn("a:tr"))[-1]))

# --- fill ---
table = tbl_shape.table
for r in range(need_rows):
    for c in range(need_cols):
        cell = table.cell(r, c)
        cell.text = DATA[r][c]
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(13)
                if r == 0:
                    run.font.bold = True

prs.save(FILE)
print(f"Table resized {cur_rows}x{cur_cols} -> {need_rows}x{need_cols} and filled.")
