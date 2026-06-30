# -*- coding: utf-8 -*-
"""Replace only the Methodology box in the filled poster with the
process-focused (less technical) version."""

from pptx import Presentation

FILE = "ERM_Poster_FILLED.pptx"

METHOD = (
    "The project follows a modular, layered design that separates the user "
    "interface, the processing/decision logic, and the data storage, allowing "
    "each part to be developed, tested, and scaled independently.\n"
    "Each emergency report is handled through a sequential AI processing "
    "pipeline: the report is received with the citizen's location and timestamp; "
    "voice reports are converted to text; the text is screened for authenticity "
    "(fake-report filtering); the incident is classified by type and severity; "
    "it is checked against recent reports for duplicates; valid incidents are "
    "placed in a priority queue; the most suitable available unit is selected and "
    "dispatched; and all parties are notified, with feedback collected after "
    "resolution.\n"
    "Dispatching follows a two-step decision method: incidents are first ordered "
    "by urgency (severity, then waiting time), and a unit is then matched to each "
    "incident by geographic proximity. When a unit becomes free, it is "
    "automatically reassigned to the highest-priority waiting incident.\n"
    "The AI models were trained on a labeled dataset of emergency descriptions "
    "and integrated into the system as an always-on processing service."
)


def first_font_size(tf):
    for p in tf.paragraphs:
        for r in p.runs:
            return r.font.size
    return None


prs = Presentation(FILE)
shape = list(prs.slides[0].shapes)[8]  # Methodology body
tf = shape.text_frame
size = first_font_size(tf)
tf.word_wrap = True
tf.clear()
for i, line in enumerate(METHOD.split("\n")):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = line
    if size is not None:
        run.font.size = size

prs.save(FILE)
print("Methodology updated in", FILE)
