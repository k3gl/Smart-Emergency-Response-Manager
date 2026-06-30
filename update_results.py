# -*- coding: utf-8 -*-
"""Reformat the Results/Findings box into the categorized style
(numbered bold headers + sub-bullets), like the reference poster."""

from pptx import Presentation
from pptx.util import Pt

FILE = "ERM_Poster_FILLED.pptx"

LEAD = "Key findings from testing the integrated system:"

# (header, [sub-bullets])
FINDINGS = [
    ("1- The AI triage pipeline performed accurately:", [
        "Severity classification (Low / Urgent / Critical) achieved high accuracy on the test set.",
        "Unit-type prediction correctly identified the required services (Police / Ambulance / Fire).",
        "OpenAI Whisper reliably transcribed voice SOS reports, even under noisy conditions.",
    ]),
    ("2- Fake and duplicate filtering improved resource efficiency:", [
        "The fake-report detector prevented unnecessary dispatches to false alarms.",
        "Duplicate detection merged repeated reports of the same event into one case.",
    ]),
    ("3- Smart dispatching reduced response time:", [
        "The nearest available unit was selected automatically using Haversine distance.",
        "Severity-based prioritization ensured critical incidents were served first.",
        "Freed units were automatically reassigned to the highest-priority waiting incident.",
    ]),
    ("4- The system demonstrated scalability and reliability:", [
        "One global queue and a centralized database coordinated all stations.",
        "Race-free dispatching prevented double-assignment of units under load.",
    ]),
    ("5- Real-time communication enhanced the user experience:", [
        "Citizens received push notifications and live tracking of the responding unit.",
        "Citizens could rate the service and leave feedback after resolution.",
    ]),
]


def first_font_size(tf):
    for p in tf.paragraphs:
        for r in p.runs:
            return r.font.size
    return None


prs = Presentation(FILE)
shapes = list(prs.slides[0].shapes)

# Short lead-in line in the top results box.
set_tf = shapes[20].text_frame
sz20 = first_font_size(set_tf)
set_tf.word_wrap = True
set_tf.clear()
r = set_tf.paragraphs[0].add_run()
r.text = LEAD
if sz20:
    r.font.size = sz20

# Categorized findings in the findings box.
tf = shapes[13].text_frame
size = first_font_size(tf) or Pt(14)
tf.word_wrap = True
tf.clear()
first = True
for header, subs in FINDINGS:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.level = 0
    run = p.add_run()
    run.text = header
    run.font.bold = True
    run.font.size = size
    for s in subs:
        sp = tf.add_paragraph()
        sp.level = 1
        sr = sp.add_run()
        sr.text = s
        sr.font.size = size

prs.save(FILE)
print("Results reformatted in", FILE)
