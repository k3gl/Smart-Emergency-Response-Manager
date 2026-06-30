# -*- coding: utf-8 -*-
"""Fill the GP poster template with the Emergency Response Manager content.
Writes a NEW file so the original template is untouched."""

from pptx import Presentation

SRC = "70x100_Vertical_Template -2026.pptx"
OUT = "ERM_Poster_FILLED.pptx"

TITLE = "Smart Emergency Response Manager"

INTRO = (
    "Traditional emergency systems depend on phone calls and manual dispatching, "
    "causing delays, fragmented information, and heavy reliance on human judgment "
    "under pressure.\n"
    "The Smart Emergency Response Manager (ERM) is an AI-driven platform connecting "
    "Citizens, Response Units, and Administrators. Citizens report emergencies by "
    "text or voice; the system automatically attaches GPS location and timestamp, "
    "then uses AI to verify, classify, prioritize, and dispatch the nearest suitable "
    "unit automatically.\n"
    "The goal is to reduce response time, automate triage, and optimize emergency "
    "resources — turning a manual, reactive process into a proactive, data-driven one."
)

METHOD = (
    "The system uses a three-layer architecture: a Flutter mobile app (Presentation) "
    "for Citizens, Units, and Admins; a Python AI microservice plus a "
    "Supabase/PostgreSQL backend (Application); and PostgreSQL with row-level security "
    "and real-time synchronization (Data).\n"
    "AI pipeline per report: voice is transcribed with OpenAI Whisper, then the report "
    "is screened for fake content, classified by severity and required unit type, and "
    "checked for duplicates. Valid incidents enter a severity-priority queue and are "
    "dispatched to the nearest available matching unit, followed by push notifications "
    "and feedback.\n"
    "The dispatch decision (priority and matching) runs inside PostgreSQL using "
    "trigger-driven functions, providing a single, race-free source of truth shared by "
    "the app and the AI worker."
)

KEYALG = (
    "• Severity: TF-IDF features + multiclass Logistic Regression (Low / Urgent / Critical).\n"
    "• Unit type: three binary Logistic Regression classifiers over 768-dim "
    "Sentence-Transformer (all-mpnet-base-v2) embeddings (Police / Ambulance / Fire).\n"
    "• Fake detection: binary classifier over the same embeddings; flagged reports "
    "are rejected before dispatch.\n"
    "• Duplicate detection: geographic (≤ 500 m), temporal (≤ 30 min), and "
    "semantic similarity (cosine ≥ 0.5).\n"
    "• Smart dispatch: nearest available unit via the Haversine great-circle "
    "distance; ETA from road-adjusted distance and average city speed.\n"
    "• Priority queue: severity order (Critical > Urgent > Low) with FIFO tie-break; "
    "a database trigger auto-reassigns freed units."
)

RESULTS_INTRO = (
    "The integrated AI pipeline processes both text and voice reports in real time and "
    "supports accurate, automated dispatch decisions:"
)

BULLETS = [
    "OpenAI Whisper transcribes voice SOS reports accurately, even under noise (sirens, traffic).",
    "AI classifiers triage incidents by severity and required unit type with high accuracy.",
    "Duplicate detection consolidates multiple reports of the same event, avoiding redundant dispatch.",
    "Fake-report filtering prevents wasted deployments to false alarms.",
    "Automated priority dispatch and trigger-based re-dispatch reduce time-to-assign.",
]

TABLE_CAPTION = "Table 1: Model Performance (test set)"

RESULTS_OUTRO = (
    "Response-time metrics (time-to-dispatch and time-to-resolve) are computed from "
    "incident timestamps to evaluate operational performance."
)

CONCLUSION = (
    "The Smart Emergency Response Manager provides a complete end-to-end intelligent "
    "workflow that existing systems lack — unifying multi-modal reporting, AI triage "
    "(fake detection, severity, unit type), duplicate handling, severity-based "
    "prioritization, geospatial smart dispatching, and real-time notifications in one "
    "platform.\n"
    "By automating these stages, it reduces emergency response times and improves "
    "resource utilization, turning a manual, reactive process into a proactive, "
    "data-driven framework.\n"
    "Future work: live road-routing for exact ETAs, a four-level severity scale, "
    "coverage-aware dispatching, and cloud deployment of the AI worker."
)

REFERENCES = (
    "1. Radford, A., et al. (2023). Robust Speech Recognition via Large-Scale Weak "
    "Supervision (Whisper). OpenAI.\n"
    "2. Reimers, N., & Gurevych, I. (2019). Sentence-BERT: Sentence Embeddings using "
    "Siamese BERT-Networks. EMNLP.\n"
    "3. Pedregosa, F., et al. (2011). Scikit-learn: Machine Learning in Python. JMLR, "
    "12, 2825–2830.\n"
    "4. Supabase Documentation. https://supabase.com/docs"
)


def first_font_size(tf):
    for p in tf.paragraphs:
        for r in p.runs:
            return r.font.size
    return None


def set_body(shape, text):
    """Replace a text box's content, preserving the original font size."""
    tf = shape.text_frame
    size = first_font_size(tf)
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if size is not None:
            run.font.size = size


def replace_in_place(shape, lines):
    """Replace paragraph text in place (keeps bullet formatting)."""
    tf = shape.text_frame
    paras = list(tf.paragraphs)
    size = first_font_size(tf)
    for i, p in enumerate(paras):
        if i < len(lines):
            if p.runs:
                p.runs[0].text = lines[i]
                for r in p.runs[1:]:
                    r.text = ""
            else:
                p.add_run().text = lines[i]
        else:
            for r in p.runs:
                r.text = ""
    for line in lines[len(paras):]:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if size is not None:
            run.font.size = size


def set_title(shape, new_title):
    for p in shape.text_frame.paragraphs:
        txt = "".join(r.text for r in p.runs)
        if "Graduation Project Title" in txt and p.runs:
            p.runs[0].text = new_title
            for r in p.runs[1:]:
                r.text = ""


prs = Presentation(SRC)
shapes = list(prs.slides[0].shapes)

set_title(shapes[6], TITLE)        # Title block
set_body(shapes[2], INTRO)         # Introduction body
set_body(shapes[8], METHOD)        # Methodology body
set_body(shapes[19], KEYALG)       # Key Algorithms body
set_body(shapes[20], RESULTS_INTRO)  # Results body (intro)
replace_in_place(shapes[13], BULLETS)  # Results bullets
set_body(shapes[16], TABLE_CAPTION)  # Table caption
set_body(shapes[21], RESULTS_OUTRO)  # Results body (after table)
set_body(shapes[10], CONCLUSION)   # Conclusion body
set_body(shapes[9], REFERENCES)    # References body

prs.save(OUT)
print("Saved:", OUT)
