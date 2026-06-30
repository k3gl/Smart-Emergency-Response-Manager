# -*- coding: utf-8 -*-
"""Replace only the Introduction box in the already-filled poster with a
longer version. Other boxes / your edits are left untouched."""

from pptx import Presentation

FILE = "ERM_Poster_FILLED.pptx"

INTRO = (
    "Emergency response is a time-critical public-safety task where every second "
    "directly affects survival and damage. Traditional systems (Public Safety "
    "Answering Points) still rely on phone calls and manual coordination between "
    "callers, dispatchers, and response units. Under heavy load or large-scale "
    "events this leads to delayed responses, fragmented and duplicated information, "
    "miscommunication, and inconsistent decisions driven by human fatigue and "
    "stress.\n"
    "The rise of Smart City and Next-Generation 911 initiatives has pushed emergency "
    "management toward digital, AI-assisted platforms that combine mobile "
    "applications, cloud services, geospatial analysis, and Artificial Intelligence "
    "to speed up decisions and improve situational awareness.\n"
    "The Smart Emergency Response Manager (ERM) follows this direction as an "
    "AI-driven platform connecting Citizens, Response Units, and Administrators in a "
    "single ecosystem. Citizens report emergencies by text or voice, and the system "
    "automatically attaches their GPS location and timestamp. The AI then transcribes "
    "voice reports, filters out fake reports, classifies the incident type and "
    "severity, and detects duplicate reports of the same event. Valid incidents are "
    "prioritized and the nearest suitable unit is dispatched automatically, with "
    "real-time notifications and live tracking for the citizen.\n"
    "The goal is to reduce response time, automate triage, optimize the use of "
    "emergency resources, and minimize reliance on manual judgment — transforming a "
    "largely manual and reactive process into a proactive, data-driven response "
    "framework."
)


def first_font_size(tf):
    for p in tf.paragraphs:
        for r in p.runs:
            return r.font.size
    return None


prs = Presentation(FILE)
shape = list(prs.slides[0].shapes)[2]  # Introduction body
tf = shape.text_frame
size = first_font_size(tf)
tf.word_wrap = True
tf.clear()
for i, line in enumerate(INTRO.split("\n")):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = line
    if size is not None:
        run.font.size = size

prs.save(FILE)
print("Introduction updated in", FILE)
